import base64
import io
import logging
import os
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import PyPDF2
from django.conf import settings
from openai import OpenAI

try:
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

try:
    import pandas as pd
except ImportError:  # pragma: no cover
    pd = None

from jobs.models import Job

logger = logging.getLogger(__name__)

SUMMARY_PROMPT = """
Attachedment Read Very care fully and All instrcution and all informtion read care fully step by step in details.
You are an AI assistant specialized in understanding writing tasks and producing a structured Job Summary, not the full content itself. Read the user's instructions and any extracted text from attachments (e.g., PDFs, DOCX) to identify what needs to be written, including topic, word count or length, reference style (APA, MLA, Harvard, etc.), and writing style or document type (essay, report, PPT, proposal, article, dissertation, thesis, etc.). If a detail is not explicitly given but can be reasonably inferred, infer it; if it cannot be inferred confidently, mark it as "Not specified." Always respond in this exact format, each on its own line and using a hyphen after the label: Topic - <short topic or title>; Word Count - <number of words or If word count is not mentioned in the Job card, then by default print "1500">; Referencing Style - <style or If Reference Style is not mentioned in the Job card, then by default print "Harvard">; Academic Style - <type or "Report">; Academic Level - Undergraduate/Masters/PhD; Summary - <What needs to be written>; Marking Criteria - <Assessment requirements>; Merit Criteria - <Excellence indicators>; Subject Field - <Discipline/area of study>; Job Summary - <10-20 sentences clearly describing what needs to be written, the main themes to cover, target audience or level if known, and any important constraints such as tone or structure>. Do not add extra sections, do not explain your reasoning, and do not write the actual assignment-only provide a clear, concise, implementation-ready Job Summary that another writer or AI could directly follow.The Job Summary must stay concise, target around 200-250 words, and never exceed 250 words. If the instructions specify a word count or range, use that value; only default to 1500 when nothing is provided.
""".strip()

STRUCTURE_PROMPT = """
You are an AI assistant specialized in creating academic writing structures (detailed outlines) for writing tasks. Your input is always the full output of a Job Summary agent, which includes at least: Topic, Word Count, Reference Style, Writing Style, and Job Summary (and may also include extra instructions). Your job is to design a clear, logically ordered, academically appropriate structure with word counts for each section and subsection, so that another writer or AI could directly draft the final document. Strictly follow all instructions and requirements from the Job Summary and ensure that every key theme, focus area, or constraint is reflected in the structure. Use academic writing conventions that match the Writing Style (e.g., essays with introduction/body/conclusion; reports with sections such as introduction, methodology, analysis, conclusion; dissertations/thesis with chapters such as introduction, literature review, methodology, results, discussion, conclusion; PPTs as slide-based academic sections, etc.). Handle Word Count as follows: always use only word counts and never pages, lines, slides, or any other length unit; if a specific word count is given, treat it as the target total and allocate section word counts so they sum to approximately that total (with minor acceptable variation); if a range is given, internally pick a reasonable midpoint and allocate based on that; if the word count is described in pages or similar, internally convert to an approximate word count and output only word counts; if Word Count is "Not specified," infer a reasonable total based on the Writing Style and academic context, then allocate accordingly. Respect the Reference Style by including a final "References" or "Bibliography" section with an appropriate word count whenever references are expected for that type of task. Ensure a coherent hierarchy with numbered sections and, where useful, subsections, each with a clear academic-style heading and an explicit word count (e.g., "Section Title - X words"). Begin by stating the title (using the Topic) and the total word count, then list the sections in order. Do not write any actual content of the sections, only the structure and word counts. Do not explain your reasoning, do not add extra metadata fields, and do not mention any unit other than words. Sub points must show word counts and the sum of sub point word counts must match the parent section total; the sum of all main sections must equal the total word count. If any subsection has its own child subsections, their word counts must sum exactly to that subsection total.
You must allocate Introduction and Conclusion to approximately 10% each of the total word count (within +/-2%). Keep other sections proportional to the remaining words. Do not use bold/markdown emphasis in headings—plain text only. If Cover Page / AI Disclaimer / References are present, keep them outside the total word count (do not allocate main total to them). Ensure that subsection word counts sum exactly to their parent, and main section totals sum exactly to the overall total.Look, Total Words count is Sum of all Main Section, and then Main Section is Sum of Sub Section, and Sub Section is Sum of Sub Section of Sub Section,like
Total words count is Sum of Main Section. and Main Section look like, 1., 2., 3., 4., .......
then Main Section is sum of Sub Section and Sub Section look like, 1.1., 1.2., 1.3., 1.4., ........
and Sub Section is sum of Sub Section of Sub Section and Sub Section of Sub Section look like 1.1.1., 1.1.2., 1.1.3., 1.1.4., ........

""".strip()

CONTENT_PROMPT = """
You are an AI assistant specialized in academic content writing. Your input is the full output of a Structure-Making Agent, which includes the title, total word count, and a numbered list of sections and subsections with individual word counts, all derived from a Job Summary. Your task is to transform this structure into complete, polished content that strictly follows all instructions, rules, and constraints implied by both the structure and the underlying task (topic, writing style, level, focus areas, tone, etc.). You must: (1) preserve the given headings and their order exactly as provided; (2) write cohesive, formal, academic prose under each section/subsection that clearly addresses the intent of its heading and the overall task; (3) follow the specified word counts closely for each section and subsection, aiming to be as close as reasonably possible to the target for each one and to the overall total; (4) maintain consistency in voice, tense, and perspective as implied by the task; and (5) ensure logical flow between sections with appropriate transitions and internal coherence. Must follow the exact word count that is mentioned, but sometimes you can provide 5% More or less in the contents as word counts. Do not modify or invent new sections, do not change the title, and do not contradict any explicit requirements from the task (such as focus, scope, or audience). When writing the content, do not include any reference list, bibliography, or citations of any kind (no in-text citations, no author-year, no numbers in brackets, and no “References” section), even if the structure or task mentions a reference style; treat that aspect as handled elsewhere. Do not explain your reasoning or describe your process; output only the final written content organized under the given headings.You must obey every user-provided input and the structure word counts: each subsection must meet its stated word target (≈exact, within a minimal tolerance), each main heading must equal the sum of its subsections, and the sum of all main headings must match the total word count from the structure. Do not alter headings, order, or totals. If any length cannot be met, prefer slight underfill (not overfill) and stay within ±2% of each target. Use plain text headings only—no markdown (#, ##), no bold (**); keep the original numbering exactly as provided.
Always obey the target total word count first, keeping your final response within ±10% of the user’s specified total (if the target is T, your answer must be between 0.9T and 1.1T words), and then strictly follow the exact section and subsection structure (headings and hierarchy) provided by the user, using the same titles and not adding any new sections. If the user also gives approximate word counts per section, treat those as guidelines while ensuring the total word count stays within the allowed range. Be concise, avoid repetition, and prioritize clarity and relevance when space is limited instead of adding extra detail. Do not mention word counts, calculations, rules, or reasoning in your output, and do not restate or reference these instructions. Your entire response should only consist of the content requested by the user, formatted using the exact structure they provided, while keeping the total word count strictly within the ±10% range.

""".strip()

REFERENCES_PROMPT = """
You are an AI assistant specialized in generating academic reference lists and corresponding in-text citation formats. Your input will be: (1) the full content produced by a content-creation agent, (2) the specified reference style (e.g., APA, MLA, Chicago, Harvard, IEEE, etc.), and (3) the approximate total word count of the content. Your task is to create an original, topic-related reference list that strictly follows the given reference style and is based on the themes, concepts, and topics present in the content. All references you provide must be to real, credible, and verifiable sources published after 2021 (i.e., from 2022 onwards). For every 1000 words of content, generate approximately 7 references (rounding reasonably to the nearest whole number) and ensure that all references are directly relevant to the subject matter of the content. Present the references as a properly formatted “Reference List” ordered alphabetically (A–Z) by the first author’s surname, strictly conforming to the rules of the specified reference style. After the alphabetical reference list, provide a separate “Citation List” that contains the in-text citation format for each reference above (e.g., for Harvard and APA: Author, Year; for MLA: Author page; for IEEE: [number], etc.), covering all references already listed. In-text Citation rules: For Harvard, APA, APA7,  IEEE Referencing (If one, two, or three authors are present in the Reference, then use the Surname of Each Author first, then a comma, and then the year in a Single bracket). Like example: ‘Hermes, A. and Riedl, R., 2021, July. Dimensions of retail customer experience and its outcomes: a literature review and directions for future research. If you notice here, two authors are present, so the in-text citation will be “(Hermes and Riedl, 2021)”. If 4 or more authors are present, then use the first author's surname, then et al., then a comma, and then the year. For example: “Pappas, A., Fumagalli, E., Rouziou, M. and Bolander, W., 2023. More than machines: The role of the future retail salesperson in enhancing the customer experience. Journal of Retailing, 99(4), pp.518-531.”. If you notice here 4 authors are present, so the intext citation will be (Pappas et al. 2023). In IEEE, all are the same but in Number Format like [1], [2], etc. Do not include any explanation, analysis, or extra text beyond the reference list and the citation list. Do not rewrite or summarize the original content. Your entire output must consist only of the formatted reference list followed by the citation list.
""".strip()

FINALIZE_PROMPT = """
You are an AI assistant specialized in finalizing academic documents by inserting in-text citations and appending an existing reference list. Your inputs are: (1) a complete piece of content with no citations or reference list, (2) a formatted reference list, (3) a citation list that specifies the correct in-text citation format for each reference, and (4) the reference style to follow (e.g., APA, MLA, Chicago, Harvard, IEEE, etc.). Your task is to cite all existing references from the citation list within the content and then append the full reference list at the end of the document, strictly following the given reference style. You must not rewrite, expand, shorten, reorder, or otherwise change any of the existing content, headings, or wording; you may only insert in-text citations at appropriate locations and add the reference list at the end. Don't cite in the Introduction, Conclusion parts, and if available, Abstract and Executive summary; in those parts, don't add in-text citations. Do not add new references, do not remove any existing references, and do not invent sources. Ensure that every reference from the provided reference list is cited at least once in the body using the corresponding in-text format from the citation list, and that all in-text citations match entries in the reference list. Maintain the original structure and formatting of the content as much as possible, only adding the necessary citation markers and the final reference list section. As output, return the full content with the in-text citations properly inserted and the complete reference list appended at the end, and do not include any explanations, notes, or extra commentary.You must not write any new content, delete any existing content, rewrite, paraphrase, reorder, correct, or explain anything in the text. You may only insert citation markers in suitable positions without changing the existing words, sentences, structure, or formatting of the content. Do not add new references that are not in the user’s reference list.
""".strip()

DEFAULT_JOB_INSTRUCTIONS = (
    "The instructions for the writing task are in the following extracted file contents. "
    "Please infer all possible details about the assignment."
)

_openai_client: Optional[OpenAI] = None


def _get_openai_client() -> OpenAI:
    """
    Lazily instantiate an OpenAI client using either the Django setting or the env var.
    """
    global _openai_client
    if _openai_client is not None:
        return _openai_client

    api_key = getattr(settings, "OPENAI_API_KEY", None) or os.environ.get("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY is not configured.")
    _openai_client = OpenAI(api_key=api_key)
    return _openai_client


def _model(setting_name: str, fallback: str) -> str:
    return getattr(settings, setting_name, fallback)


def extract_text_from_pdf(file_path: str) -> str:
    """Best-effort extraction from PDF."""
    text_parts: List[str] = []
    try:
        with open(file_path, "rb") as handle:
            reader = PyPDF2.PdfReader(handle)
            for page in reader.pages:
                snippet = page.extract_text() or ""
                if snippet:
                    text_parts.append(snippet)
    except Exception as exc:  # pragma: no cover - defensive logging
        logger.warning("Failed to read PDF %s: %s", file_path, exc)
    return "\n\n".join(text_parts).strip()


def extract_text_from_docx(file_path: str) -> str:
    """Best-effort extraction from DOCX."""
    if Document is None:
        return "[python-docx is not installed on this server.]"
    try:
        doc = Document(file_path)
        return "\n".join(para.text for para in doc.paragraphs).strip()
    except Exception as exc:  # pragma: no cover - defensive logging
        logger.warning("Failed to read DOCX %s: %s", file_path, exc)
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    if Presentation is None:
        return "[python-pptx is not installed on this server.]"
    texts: List[str] = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
    except Exception as exc:  # pragma: no cover
        logger.warning("Failed to read PPTX %s: %s", file_path, exc)
    return "\n".join(texts).strip()


def extract_text_from_csv(file_path: str) -> str:
    if pd is None:
        return "[pandas is required to parse CSV files.]"
    try:
        df = pd.read_csv(file_path)
        return df.to_csv(index=False)
    except Exception as exc:  # pragma: no cover
        logger.warning("Failed to read CSV %s: %s", file_path, exc)
        return ""


def extract_text_from_excel(file_path: str) -> str:
    if pd is None:
        return "[pandas with openpyxl is required to parse Excel files.]"
    try:
        df = pd.read_excel(file_path)
        return df.to_csv(index=False)
    except Exception as exc:  # pragma: no cover
        logger.warning("Failed to read Excel %s: %s", file_path, exc)
        return ""


def extract_text_from_plain(file_path: str) -> str:
    try:
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")
    except Exception as exc:  # pragma: no cover
        logger.warning("Failed to read text file %s: %s", file_path, exc)
        return ""


def extract_text_from_attachments(job: Job) -> Tuple[str, List[Dict[str, str]]]:
    """
    Load and return a tuple of (combined_text_blocks, image_inputs) for all attachments.
    Image inputs follow the OpenAI vision format (base64 data URLs).
    """
    text_blocks: List[str] = []
    image_contents: List[Dict[str, str]] = []

    for attachment in job.attachments.all():
        file_path = attachment.file.path
        ext = os.path.splitext(attachment.filename.lower())[1]

        if ext in {".png", ".jpg", ".jpeg"}:
            try:
                raw = Path(file_path).read_bytes()
                mime = "image/png" if ext == ".png" else "image/jpeg"
                b64 = base64.b64encode(raw).decode("utf-8")
                image_contents.append(
                    {"type": "input_image", "image_url": f"data:{mime};base64,{b64}"}
                )
            except OSError as exc:  # pragma: no cover
                logger.warning("Failed to read image %s: %s", file_path, exc)
            continue

        if ext == ".pdf":
            text = extract_text_from_pdf(file_path)
        elif ext in {".docx"}:
            text = extract_text_from_docx(file_path)
        elif ext == ".doc":
            # Best-effort: try plain text read for legacy .doc
            text = extract_text_from_plain(file_path)
        elif ext in {".pptx"}:
            text = extract_text_from_pptx(file_path)
        elif ext == ".csv":
            text = extract_text_from_csv(file_path)
        elif ext in {".xlsx", ".xls", ".xlx"}:
            text = extract_text_from_excel(file_path)
        else:
            text = extract_text_from_plain(file_path)

        # Fallback raw decode if empty
        if not text:
            try:
                raw = Path(file_path).read_bytes()
                text = raw.decode("utf-8", errors="ignore")
            except Exception:
                text = ""

        if text:
            text_blocks.append(f"----- File: {attachment.filename} -----\n{text}")

    combined_text = "\n\n".join(text_blocks).strip()
    return combined_text, image_contents


def _run_openai_request(
    *,
    instructions: str,
    content_items: List[Dict[str, str]],
    model: str,
) -> Optional[str]:
    try:
        client = _get_openai_client()
        messages = [{"role": "user", "content": content_items}]
        response = client.responses.create(
            model=model,
            instructions=instructions,
            input=messages,
        )
        output_text = getattr(response, "output_text", None)
        if not output_text and hasattr(response, "output"):
            # Fallback for SDK versions that return the chunks
            output_text = "".join(getattr(chunk, "content", [""])[0].text for chunk in response.output)
        return (output_text or "").strip()
    except Exception as exc:  # pragma: no cover - depends on external API
        logger.exception("OpenAI request failed: %s", exc)
        return None


def parse_job_summary(summary_text: str) -> Dict[str, object]:
    """
    Convert the agent output into structured data.
    """
    data: Dict[str, object] = {
        "topic": "Not specified",
        "word_count": 1500,
        "reference_style": "Harvard",
        "writing_style": "Report",
        "summary": summary_text.strip(),
    }

    parts = summary_text.split(";")
    summary_started = False
    additional_summary: List[str] = []

    for part in parts:
        segment = part.strip()
        if not segment or ":" not in segment:
            if summary_started and segment:
                additional_summary.append(segment)
            continue

        key, value = segment.split(":", 1)
        key_lower = key.strip().lower()
        value = value.strip()

        if key_lower == "topic":
            data["topic"] = value or data["topic"]
        elif key_lower == "word count":
            digits = "".join(ch for ch in value if ch.isdigit())
            try:
                data["word_count"] = int(digits) if digits else data["word_count"]
            except ValueError:
                pass
        elif key_lower == "reference style":
            data["reference_style"] = value or data["reference_style"]
        elif key_lower == "writing style":
            data["writing_style"] = value or data["writing_style"]
        elif key_lower == "job summary":
            data["summary"] = value
            summary_started = True
        elif summary_started:
            additional_summary.append(segment)

    if summary_started and additional_summary:
        data["summary"] = f"{data['summary']}; {'; '.join(additional_summary)}"

    return data


def _extract_word_count_hint(text: str) -> Optional[str]:
    """
    Find word count or range hints in free text.
    """
    text = text or ""
    # normalize commas in numbers, e.g., 2,500 -> 2500
    text = re.sub(r"(?<=\d),(?=\d)", "", text)
    patterns = [
        r"(\d{2,5})\s*-\s*(\d{2,5})\s*words?",
        r"(\d{2,5})\s*[–—\-to]{1,3}\s*(\d{2,5})\s*words?",
        r"(\d{2,5})\s*words?\b",
        r"words?\s*[:\-]?\s*(\d{2,5})",
        r"word\s*count\s*[:\-]?\s*(\d{2,5})",
        r"\bwc\s*[:\-]?\s*(\d{2,5})\b",
        r"word\s*limit\s*(?:of\s*)?[:\-]?\s*(\d{2,5})(?:\s*[–—\-to]{1,3}\s*(\d{2,5}))?",
    ]
    for pat in patterns:
        for m in re.finditer(pat, text, flags=re.IGNORECASE):
            if m.lastindex and m.lastindex >= 2 and m.group(2):
                return f"{m.group(1)}-{m.group(2)}"
            if m.lastindex and m.group(1):
                return m.group(1)
    # page hints -> words (~275/page)
    page_pat = re.compile(r"(\d{1,3})(?:\s*[–—\-to]{1,3}\s*(\d{1,3}))?\s*pages?", re.IGNORECASE)
    for m in page_pat.finditer(text):
        start = int(m.group(1))
        end = int(m.group(2)) if m.group(2) else None
        if end:
            return f"{start*275}-{end*275}"
        return str(start * 275)
    return None


def _extract_ref_style(text: str) -> Optional[str]:
    """
    Detect referencing style from free text.
    """
    style_map = {
        "APA7": "APA",
        "APA 7": "APA",
        "APA 7TH": "APA",
        "APA": "APA",
        "MLA": "MLA",
        "HARVARD": "Harvard",
        "CHICAGO": "CHICAGO",
        "IEEE": "IEEE",
        "VANCOUVER": "VANCOUVER",
        "OSCOLA": "OSCOLA",
        "TURABIAN": "TURABIAN",
        "REFERENCING STYLE APA": "APA",
        "REFERENCING STYLE HARVARD": "Harvard",
        "APA SYSTEM": "APA",
        "Ieee": "IEEE",
    }
    upper_text = (text or "").upper()
    for key, val in style_map.items():
        if key in upper_text:
            return val
    return None


def _rebalance_structure_text(struct_text: str, expected_total: Optional[int] = None) -> str:
    """
    Post-process a generated structure to enforce:
    - No markdown bold (**)
    - Cover / AI Disclaimer / References excluded from totals
    - Subsection sums equal parent
    - Main sections sum to expected_total (if provided) or natural sum
    - Introduction and Conclusion near 10% each (±2%) of total
    """
    if not struct_text:
        return struct_text

    # strip bold markers
    struct_text = re.sub(r"\*\*(.*?)\*\*", r"\1", struct_text)

    ignore_keys = ['cover', 'cover page', 'ai disclaimer', 'disclaimer', 'references', 'reference', 'bibliography']

    def _is_ignored(line: str) -> bool:
        low = (line or '').lower()
        return any(k in low for k in ignore_keys)

    def _find_count(line: str):
        m = re.search(r"(\d{1,6})\s*words?", line, flags=re.IGNORECASE)
        return int(m.group(1)) if m else None

    def _replace_count(line: str, new_val: int) -> str:
        return re.sub(r"(\d{1,6})(\s*words?)", fr"{new_val}\2", line, count=1, flags=re.IGNORECASE)

    try:
        target_total = int(expected_total) if expected_total not in (None, '', 'Not specified') else None
    except Exception:
        target_total = None

    lines = struct_text.splitlines()
    total_idx = next((i for i, ln in enumerate(lines) if 'total word count' in ln.lower()), None)

    heading_re = re.compile(r'^\s*(\d+)\.\s')
    subheading_re = re.compile(r'^\s*(\d+)\.(\d+)\s')

    mains = []   # (idx, num, count, line_text)
    subs = {}    # num -> list[(idx, subnum, count)]

    for idx, line in enumerate(lines):
        if idx == total_idx:
            continue
        count = _find_count(line)
        if count is None or _is_ignored(line):
            continue
        m_sub = subheading_re.match(line)
        if m_sub:
            pnum = int(m_sub.group(1))
            subs.setdefault(pnum, []).append((idx, int(m_sub.group(2)), count))
            continue
        m_main = heading_re.match(line)
        if m_main:
            mains.append((idx, int(m_main.group(1)), count, line))

    if not mains:
        return struct_text

    # Parent counts based on subs when available
    main_counts = {}
    for idx, num, count, _line in mains:
        if num in subs:
            subtotal = sum(c[2] for c in subs[num])
            main_counts[num] = subtotal if subtotal > 0 else count
        else:
            main_counts[num] = count

    main_sum = sum(main_counts.values())
    if main_sum <= 0:
        return struct_text

    # Rescale mains to target_total if provided
    if target_total:
        scale = target_total / main_sum
        scaled = {k: max(1, round(v * scale)) for k, v in main_counts.items()}
        drift = target_total - sum(scaled.values())
        if drift != 0:
            largest = max(scaled, key=scaled.get)
            scaled[largest] = max(1, scaled[largest] + drift)
        main_counts = scaled

    # Enforce intro/conclusion ~10% each if present
    def _find_main_by_name(name_fragment: str):
        for idx, num, _count, line in mains:
            if name_fragment in line.lower():
                return num
        return None

    intro_num = _find_main_by_name('introduction')
    concl_num = _find_main_by_name('conclusion')
    total_main = sum(main_counts.values())
    target_intro = target_concl = None
    if total_main > 0:
        target_intro = round(total_main * 0.10)
        target_concl = round(total_main * 0.10)
        wiggle = round(total_main * 0.02)
        if intro_num and intro_num in main_counts:
            main_counts[intro_num] = max(1, target_intro)
        if concl_num and concl_num in main_counts:
            main_counts[concl_num] = max(1, target_concl)
        # rescale remaining to keep total
        remaining_nums = [n for n in main_counts.keys() if n not in (intro_num, concl_num)]
        remaining_sum = sum(main_counts[n] for n in remaining_nums)
        if remaining_nums and total_main > 0:
            new_total = sum(main_counts.values())
            drift = total_main - new_total
            if drift != 0:
                # push drift into largest remaining or intro if none
                adjust_target = remaining_nums[0]
                adjust_target = max(remaining_nums, key=lambda n: main_counts[n]) if remaining_nums else intro_num or concl_num
                if adjust_target:
                    main_counts[adjust_target] = max(1, main_counts[adjust_target] + drift)

    # Rescale subs to their parent
    for pnum, items in subs.items():
        parent_target = main_counts.get(pnum)
        orig = sum(c[2] for c in items)
        if not parent_target or orig <= 0:
            continue
        factor = parent_target / orig
        new_counts = [max(1, round(c[2] * factor)) for c in items]
        drift = parent_target - sum(new_counts)
        if drift != 0:
            adjust_idx = max(range(len(new_counts)), key=lambda i: new_counts[i])
            new_counts[adjust_idx] = max(1, new_counts[adjust_idx] + drift)
        for (idx_line, _, _), new_val in zip(items, new_counts):
            lines[idx_line] = _replace_count(lines[idx_line], new_val)

    # Write main counts
    for idx_line, num, _count, _line in mains:
        new_val = main_counts.get(num)
        if new_val is not None:
            lines[idx_line] = _replace_count(lines[idx_line], new_val)

    final_total = sum(main_counts.values())
    if total_idx is not None:
        lines[total_idx] = re.sub(r'(Total\s*Word\s*Count\s*[:\-]?\s*)(\d{1,6})',
                                  fr"\1{final_total}", lines[total_idx], count=1, flags=re.IGNORECASE)
    else:
        lines.insert(1, f"Total Word Count: {final_total}")

    return "\n".join(lines)

def generate_job_summary(job: Job) -> Tuple[Optional[Dict[str, object]], Optional[str]]:
    attachment_text, image_contents = extract_text_from_attachments(job)

    base_instruction = (job.instruction or "").strip()
    if not base_instruction:
        base_instruction = DEFAULT_JOB_INSTRUCTIONS

    if attachment_text:
        combined_text = (
            f"{base_instruction}\n\nBelow is the extracted text from the uploaded files:\n\n"
            f"{attachment_text}"
        )
    else:
        combined_text = base_instruction

    wc_hint = _extract_word_count_hint(combined_text)
    ref_hint = _extract_ref_style(combined_text)

    content_items: List[Dict[str, str]] = [{"type": "input_text", "text": combined_text}]
    content_items.extend(image_contents)

    model = _model("OPENAI_MODEL_SUMMARY", "gpt-4.1-mini")
    response_text = _run_openai_request(
        instructions=SUMMARY_PROMPT,
        content_items=content_items,
        model=model,
    )

    if not response_text:
        return None, "Failed to generate Job Summary."

    parsed = parse_job_summary(response_text)
    # Override defaults with detected hints from attachments/instructions
    if wc_hint:
        try:
            parsed["word_count"] = int(str(wc_hint).split("-")[0])
        except Exception:
            pass
    if ref_hint:
        parsed["reference_style"] = ref_hint
    return parsed, None


def generate_job_structure(
    job_summary_text: str,
) -> Tuple[Optional[str], Optional[str]]:
    if not job_summary_text:
        return None, "Job Summary text is required."

    wc_hint = _extract_word_count_hint(job_summary_text)
    target_wc = None
    if wc_hint:
        try:
            target_wc = int(str(wc_hint).split("-")[0])
        except Exception:
            target_wc = None

    model = _model("OPENAI_MODEL_STRUCTURE", "gpt-4.1-mini")
    response_text = _run_openai_request(
        instructions=STRUCTURE_PROMPT,
        content_items=[{"type": "input_text", "text": job_summary_text}],
        model=model,
    )

    if not response_text:
        return None, "Failed to generate Job Structure."

    cleaned = _rebalance_structure_text(response_text.strip(), expected_total=target_wc)
    return cleaned, None


def generate_content(job_structure_text: str) -> Tuple[Optional[str], Optional[str]]:
    if not job_structure_text:
        return None, "Job Structure text is required."

    def _parse_total_from_structure(struct_text: str) -> Optional[int]:
        m = re.search(r"Total\s*Word\s*Count\s*[:\-]?\s*(\d{2,6})", struct_text, flags=re.IGNORECASE)
        if m:
            try:
                return int(m.group(1))
            except Exception:
                return None
        # fallback: sum top-level headings
        total = 0
        found = False
        for ln in struct_text.splitlines():
            m2 = re.search(r"^\s*\d+\.\s.*?(\d{2,6})\s*words?", ln, flags=re.IGNORECASE)
            if m2:
                found = True
                try:
                    total += int(m2.group(1))
                except Exception:
                    pass
        return total if found else None

    def _strip_markdown(text: str) -> str:
        if not text:
            return text
        lines = []
        for line in text.splitlines():
            line = re.sub(r'^\s*#+\s*', '', line)  # remove leading #/## etc.
            line = line.replace('**', '')          # strip bold markers
            lines.append(line)
        return "\n".join(lines)

    def _strip_references_section(text: str) -> str:
        """
        Remove any trailing References block that might have been added erroneously.
        """
        if not text:
            return text
        parts = re.split(r'\n\s*references\s*\n', text, flags=re.IGNORECASE)
        return parts[0].strip() if parts else text

    target_total = _parse_total_from_structure(job_structure_text)
    model = _model("OPENAI_MODEL_CONTENT", "gpt-4.1-mini")

    def _call(items: List[Dict[str, str]]) -> Optional[str]:
        return _run_openai_request(
            instructions=CONTENT_PROMPT,
            content_items=items,
            model=model,
        )

    response_text = _call([{"type": "input_text", "text": job_structure_text}])
    if not response_text:
        return None, "Failed to generate academic content."

    cleaned = _strip_references_section(_strip_markdown(response_text.strip()))
    needs_retry = False
    hash_present = '#' in cleaned
    if target_total:
        word_count = len(re.findall(r"\w+", cleaned))
        if word_count < 0.98 * target_total or word_count > 1.02 * target_total:
            needs_retry = True
    if hash_present:
        needs_retry = True
    if needs_retry:
        reminder = (
            f"Regenerate in plain text headings (no # or **). "
            f"Match every section/subsection word target exactly and keep total around {target_total or 'the stated total'} words."
        )
        second = _call([
            {"type": "input_text", "text": job_structure_text},
            {"type": "input_text", "text": reminder},
        ])
        if second:
            cleaned = _strip_references_section(_strip_markdown(second.strip()))


    return cleaned, None


def parse_references_block(raw_text: str) -> Tuple[str, str]:
    """
    Split the agent output into Reference List and Citation List sections.
    """
    if "Citation List" in raw_text:
        head, tail = raw_text.split("Citation List", 1)
        return head.strip(), f"Citation List{tail.strip()}"
    return raw_text.strip(), ""


def generate_references(
    content_text: str,
    reference_style: str,
    total_words: int,
) -> Tuple[Optional[Dict[str, str]], Optional[str]]:
    if not content_text:
        return None, "Content text is required to build references."

    combined = (
        f"Reference style: {reference_style or 'Harvard'}\n"
        f"Approximate total word count: {total_words or 0}\n\n"
        f"Content:\n{content_text}"
    )

    model = _model("OPENAI_MODEL_REFERENCES", "gpt-4.1")
    response_text = _run_openai_request(
        instructions=REFERENCES_PROMPT,
        content_items=[{"type": "input_text", "text": combined}],
        model=model,
    )

    if not response_text:
        return None, "Failed to generate references."

    reference_list, citation_list = parse_references_block(response_text)
    return {"reference_list": reference_list, "citation_list": citation_list}, None


def generate_full_content_with_citations(
    content_text: str,
    reference_list: str,
    citation_list: str,
    reference_style: str,
) -> Tuple[Optional[str], Optional[str]]:
    if not content_text or not reference_list:
        return None, "Content and reference list are required."

    combined = (
        f"Reference style: {reference_style or 'Harvard'}\n\n"
        "=== CONTENT (NO CITATIONS) ===\n"
        f"{content_text}\n\n"
        "=== REFERENCE LIST ===\n"
        f"{reference_list}\n\n"
        "=== CITATION LIST ===\n"
        f"{citation_list or ''}\n"
    )

    model = _model("OPENAI_MODEL_FINAL", "gpt-4.1")
    response_text = _run_openai_request(
        instructions=FINALIZE_PROMPT,
        content_items=[{"type": "input_text", "text": combined}],
        model=model,
    )

    if not response_text:
        return None, "Failed to generate final cited document."

    return response_text.strip(), None


def check_plagiarism(content_text: str) -> Tuple[Optional[Dict[str, object]], Optional[str]]:
    if not content_text:
        return None, "Content text is required for plagiarism check."

    similarity = 5.2
    report = (
        "Plagiarism Report\n\n"
        "This is a placeholder report. Integrate with a real plagiarism service.\n"
        f"Total Word Count: {len(content_text.split())}\n"
        f"Similarity Percentage: {similarity}%\n"
        "Status: PASSED\n"
    )
    return {"report": report, "similarity_percentage": similarity}, None


def check_ai_content(content_text: str) -> Tuple[Optional[Dict[str, object]], Optional[str]]:
    if not content_text:
        return None, "Content text is required for AI detection check."

    ai_percentage = 12.5
    report = (
        "AI Detection Report\n\n"
        "This is a placeholder report. Integrate with a real AI-content detection API.\n"
        f"Total Word Count: {len(content_text.split())}\n"
        f"AI-Generated Content: {ai_percentage}%\n"
        "Status: ACCEPTABLE\n"
    )
    return {"report": report, "ai_percentage": ai_percentage}, None


